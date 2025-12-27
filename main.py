# main.py
import shutil
import subprocess
import tempfile
from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse
from typing import List, Optional
from pydantic import BaseModel
import os
import uuid
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from pdf2docx import Converter
import camelot
import pytesseract
from PIL import Image
import io
import pandas as pd

app = FastAPI(title="PDF Tools API")

# Helper: save UploadFile to a temp file and return path
def save_uploadfile_tmp(upload_file: UploadFile) -> str:
    suffix = os.path.splitext(upload_file.filename)[1] or ""
    fd, path = tempfile.mkstemp(suffix=suffix)
    os.close(fd)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(upload_file.file, buffer)
    return path

# HEALTH (use this for external pings)
@app.get("/ping")
async def ping():
    return {"status": "ok", "message": "pong"}

# 1. PDF -> Word (pdf2docx)
@app.post("/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    out = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    out_path = out.name
    out.close()
    try:
        cv = Converter(path)
        cv.convert(out_path, start=0, end=None)
        cv.close()
        return StreamingResponse(open(out_path, "rb"), media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                 headers={"Content-Disposition": f"attachment; filename={os.path.basename(out_path)}"})
    finally:
        try: os.remove(path)
        except: pass

# 2. Word -> PDF (libreoffice headless)
@app.post("/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    tmp_dir = tempfile.mkdtemp()
    try:
        # libreoffice will write into cwd or specified outdir
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, path], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        base = os.path.splitext(os.path.basename(path))[0] + ".pdf"
        pdf_path = os.path.join(tmp_dir, base)
        if not os.path.exists(pdf_path):
            raise HTTPException(status_code=500, detail="Conversion failed")
        return StreamingResponse(open(pdf_path, "rb"), media_type="application/pdf",
                                 headers={"Content-Disposition": f"attachment; filename={base}"})
    finally:
        try: os.remove(path)
        except: pass

# 3. PDF -> JPG (export each page as JPG, return zip if multiple)
@app.post("/convert/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    images = convert_from_path(path, dpi=200)
    # If single page return single jpg, else create zip
    if len(images) == 1:
        buf = io.BytesIO()
        images[0].save(buf, format="JPEG")
        buf.seek(0)
        return StreamingResponse(buf, media_type="image/jpeg", headers={"Content-Disposition": "attachment; filename=page1.jpg"})
    else:
        import zipfile
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, mode="w") as zf:
            for i, img in enumerate(images, start=1):
                b = io.BytesIO()
                img.save(b, format="JPEG")
                b.seek(0)
                zf.writestr(f"page_{i}.jpg", b.read())
        zip_buf.seek(0)
        return StreamingResponse(zip_buf, media_type="application/zip", headers={"Content-Disposition": "attachment; filename=pages.zip"})
    # cleanup
    try: os.remove(path)
    except: pass

# 4. JPG -> PDF
@app.post("/convert/jpg-to-pdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...)):
    images = []
    paths = []
    for f in files:
        p = save_uploadfile_tmp(f)
        paths.append(p)
        images.append(Image.open(p).convert("RGB"))
    out_path = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    images[0].save(out_path, save_all=True, append_images=images[1:])
    for p in paths:
        try: os.remove(p)
        except: pass
    return StreamingResponse(open(out_path, "rb"), media_type="application/pdf", headers={"Content-Disposition": "attachment; filename=converted.pdf"})

# 5. PDF -> Excel (Camelot)
@app.post("/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    try:
        # Camelot: 'lattice' or 'stream' detection; try both
        tables = []
        try:
            tables = camelot.read_pdf(path, pages='all', flavor='lattice')
            if len(tables) == 0:
                tables = camelot.read_pdf(path, pages='all', flavor='stream')
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Camelot error: {e}")
        if len(tables) == 0:
            raise HTTPException(status_code=404, detail="No tables found")
        out = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        out_path = out.name
        out.close()
        with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
            for i, t in enumerate(tables, start=1):
                df = t.df
                df.to_excel(writer, sheet_name=f"table_{i}", index=False)
        return StreamingResponse(open(out_path, "rb"), media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                 headers={"Content-Disposition": f"attachment; filename={os.path.basename(out_path)}"})
    finally:
        try: os.remove(path)
        except: pass

# 6. Excel -> PDF (libreoffice)
@app.post("/convert/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    tmp_dir = tempfile.mkdtemp()
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, path], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        base = os.path.splitext(os.path.basename(path))[0] + ".pdf"
        pdf_path = os.path.join(tmp_dir, base)
        if not os.path.exists(pdf_path):
            raise HTTPException(status_code=500, detail="Conversion failed")
        return StreamingResponse(open(pdf_path, "rb"), media_type="application/pdf",
                                 headers={"Content-Disposition": f"attachment; filename={base}"})
    finally:
        try: os.remove(path)
        except: pass

# --- Page Manipulation Tools (using PyMuPDF / pypdf or fitz) ---
@app.post("/tools/merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    out_path = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc = fitz.open()
    for f in files:
        p = save_uploadfile_tmp(f)
        src = fitz.open(p)
        doc.insert_pdf(src)
        src.close()
        try: os.remove(p)
        except: pass
    doc.save(out_path)
    doc.close()
    return StreamingResponse(open(out_path, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=merged.pdf"})

@app.post("/tools/split")
async def split_pdf(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    out_files = []
    for i in range(doc.page_count):
        new = fitz.open()
        new.insert_pdf(doc, from_page=i, to_page=i)
        fpath = tempfile.NamedTemporaryFile(suffix=f"_page_{i+1}.pdf", delete=False).name
        new.save(fpath)
        new.close()
        out_files.append(fpath)
    doc.close()
    # return zip
    import zipfile, io
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, mode="w") as zf:
        for p in out_files:
            zf.write(p, arcname=os.path.basename(p))
    zip_buf.seek(0)
    # cleanup
    try: os.remove(path)
    except: pass
    for p in out_files:
        try: os.remove(p)
        except: pass
    return StreamingResponse(zip_buf, media_type="application/zip", headers={"Content-Disposition":"attachment; filename=pages.zip"})

class PagesModel(BaseModel):
    pages: List[int]

@app.post("/tools/extract")
async def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):
    # pages form: comma-separated pages e.g. "1,3,5-7"
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    sel = parse_page_ranges(pages, doc.page_count)
    new = fitz.open()
    for p in sel:
        new.insert_pdf(doc, from_page=p-1, to_page=p-1)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    new.save(out)
    new.close(); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition": "attachment; filename=extracted.pdf"})

@app.post("/tools/delete-pages")
async def delete_pages(file: UploadFile = File(...), pages: str = Form(...)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    sel = parse_page_ranges(pages, doc.page_count)
    # Remove selected pages (work from end to start)
    for p in sorted(sel, reverse=True):
        doc.delete_page(p-1)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc.save(out); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=updated.pdf"})

@app.post("/tools/reorder")
async def reorder_pages(file: UploadFile = File(...), order: str = Form(...)):
    # order e.g. "2,1,3,5,4"
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    order_list = [int(x) for x in order.split(",")]
    new = fitz.open()
    for p in order_list:
        new.insert_pdf(doc, from_page=p-1, to_page=p-1)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    new.save(out); new.close(); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=reordered.pdf"})

@app.post("/tools/rotate")
async def rotate_pages(file: UploadFile = File(...), page: int = Form(...), degrees: int = Form(...)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    page_obj = doc.load_page(page-1)
    page_obj.set_rotation(degrees)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc.save(out); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=rotated.pdf"})

# Utilities
def parse_page_ranges(s: str, max_page: int):
    parts = [p.strip() for p in s.split(",")]
    pages = []
    for part in parts:
        if "-" in part:
            a,b = part.split("-")
            pages += list(range(int(a), int(b)+1))
        else:
            pages.append(int(part))
    pages = [p for p in pages if 1 <= p <= max_page]
    return sorted(set(pages))

# 13. Add Text Watermark
@app.post("/tools/watermark-text")
async def watermark_text(file: UploadFile = File(...), text: str = Form(...), fontsize: int = Form(36)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    for page in doc:
        rect = page.rect
        page.insert_text((rect.width/4, rect.height/2), text, fontsize=fontsize, rotate=45, render_mode=3, color=(0.5,0.5,0.5))
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc.save(out); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=watermarked.pdf"})

# 14. Add Page Numbers
@app.post("/tools/add-page-numbers")
async def add_page_numbers(file: UploadFile = File(...), start: int = Form(1)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    for i, page in enumerate(doc, start=start):
        page.insert_text((page.rect.width - 50, page.rect.height - 30), str(i), fontsize=12)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc.save(out); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=with-pagenumbers.pdf"})

# 15. PDF Editing (add text) (same as edit/add-text)
@app.post("/tools/edit/add-text")
async def add_text(file: UploadFile = File(...), page: int = Form(...), x: float = Form(...), y: float = Form(...), text: str = Form(...), fontsize: int = Form(12)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    p = doc.load_page(page-1)
    p.insert_text((x,y), text, fontsize=fontsize)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc.save(out); doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=edited.pdf"})

# 16. Protect PDF (password)
@app.post("/tools/protect")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    path = save_uploadfile_tmp(file)
    doc = fitz.open(path)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    # PyMuPDF supports encryption via saveAs
    doc.save(out, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password, user_pw=password)
    doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=protected.pdf"})

# 17. Unlock PDF (remove password)
@app.post("/tools/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    path = save_uploadfile_tmp(file)
    try:
        doc = fitz.open(path, filetype="pdf", password=password)
    except Exception as e:
        raise HTTPException(status_code=401, detail="Wrong password")
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    doc.save(out)
    doc.close()
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=unlocked.pdf"})

# 18. Repair PDF (Ghostscript fix)
@app.post("/tools/repair")
async def repair_pdf(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    out_path = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    try:
        subprocess.run(["gs", "-o", out_path, "-sDEVICE=pdfwrite", "-dPDFSETTINGS=/prepress", path], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"Ghostscript failed: {e}")
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out_path, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=repaired.pdf"})

# 19. Convert to PDF/A
@app.post("/tools/pdfa")
async def pdf_to_pdfa(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    out_path = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    try:
        subprocess.run(["gs", "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-sProcessColorModel=DeviceCMYK", "-sDEVICE=pdfwrite", f"-sOutputFile={out_path}", path], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"Ghostscript failed: {e}")
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out_path, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=pdfa.pdf"})

# 20. OCR PDF -> Text
@app.post("/tools/ocr")
async def ocr_pdf(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    images = convert_from_path(path, dpi=200)
    full_text = []
    for img in images:
        txt = pytesseract.image_to_string(img)
        full_text.append(txt)
    try: os.remove(path)
    except: pass
    return JSONResponse({"text": "\n".join(full_text)})

# 21. HTML -> PDF (wkhtmltopdf or weasyprint)
@app.post("/convert/html-to-pdf")
async def html_to_pdf(html: str = Form(...)):
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    # use wkhtmltopdf if installed
    p = subprocess.run(["wkhtmltopdf", "-", out], input=html.encode("utf-8"), check=False)
    if not os.path.exists(out):
        raise HTTPException(status_code=500, detail="Conversion failed")
    return StreamingResponse(open(out, "rb"), media_type="application/pdf", headers={"Content-Disposition":"attachment; filename=out.pdf"})

# Extra tools
# 22. PDF -> PPT (image-based slides)
@app.post("/convert/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    from pptx import Presentation
    from pptx.util import Inches
    path = save_uploadfile_tmp(file)
    images = convert_from_path(path, dpi=150)
    prs = Presentation()
    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img_buf = io.BytesIO()
        img.save(img_buf, format="PNG")
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    out_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
    prs.save(out_path)
    try: os.remove(path)
    except: pass
    return StreamingResponse(open(out_path, "rb"), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                             headers={"Content-Disposition":"attachment; filename=converted.pptx"})

# 23. PPT -> PDF (libreoffice)
@app.post("/convert/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    path = save_uploadfile_tmp(file)
    tmp_dir = tempfile.mkdtemp()
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, path], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        base = os.path.splitext(os.path.basename(path))[0] + ".pdf"
        pdf_path = os.path.join(tmp_dir, base)
        if not os.path.exists(pdf_path):
            raise HTTPException(status_code=500, detail="Conversion failed")
        return StreamingResponse(open(pdf_path, "rb"), media_type="application/pdf",
                                 headers={"Content-Disposition": f"attachment; filename={base}"})
    finally:
        try: os.remove(path)
        except: pass

# simple root
@app.get("/")
def root():
    return {"message": "PDF Tools API - see /docs for interactive API"}

# Run with: uvicorn main:app --host 0.0.0.0 --port $PORT
